from flask import Flask, request, jsonify, Response
import requests
import os
from docx import Document
import PyPDF2
import json
from flask_cors import CORS

app = Flask(__name__)
CORS(app)

def read_file_content(file_path):
    if file_path.endswith('.txt'):
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    elif file_path.endswith('.docx'):
        from docx import Document
        doc = Document(file_path)
        return '\n'.join([para.text for para in doc.paragraphs])
    elif file_path.endswith('.pdf'):
        import PyPDF2
        text = ""
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text
        return text
    elif file_path.endswith('.pptx'):
        from pptx import Presentation
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return '\n'.join(text_runs)
    else:
        return ""

@app.route('/api/upload_file', methods=['POST'])
def upload_file():
    file = request.files.get('file')
    if not file:
        return jsonify({'error': 'No file uploaded'}), 400
    filename = file.filename
    save_path = os.path.join('uploads', filename)
    os.makedirs('uploads', exist_ok=True)
    file.save(save_path)
    return jsonify({'filename': filename})

@app.route('/api/preview_file', methods=['GET'])
def preview_file():
    filename = request.args.get('filename')
    if not filename:
        return jsonify({'content': ''})
    save_path = os.path.join('uploads', filename)
    if not os.path.exists(save_path):
        return jsonify({'content': ''})
    content = read_file_content(save_path)
    if not content:
        content = ""
    return jsonify({'content': content})

@app.route('/api/stream_generate', methods=['POST'])
def stream_generate():
    user_prompt = request.json.get('user_prompt', '').strip()
    filename = request.json.get('filename', '').strip()
    file_content = ""
    if filename:
        save_path = os.path.join('uploads', filename)
        if os.path.exists(save_path):
            file_content = read_file_content(save_path)
    # user_prompt 必须有内容
    if not user_prompt:
        user_prompt = "请根据以下内容生成5道单项选择题，每题4个选项，并给出正确答案。"
    prompt = user_prompt
    if file_content:
        prompt += f"\n内容如下：{file_content}"
    def generate():
        url = "http://localhost:11434/api/generate"
        payload = {
            "model": "deepseek-r1:7b",
            "prompt": prompt
        }
        response = requests.post(url, json=payload, stream=True)
        for line in response.iter_lines():
            if line:
                try:
                    data = line.decode('utf-8')
                    obj = json.loads(data)
                    if "response" in obj:
                        yield obj['response']
                except Exception as e:
                    continue
    return Response(generate(), mimetype='text/plain')

if __name__ == '__main__':
    app.run(debug=True)
